#!/usr/bin/env python3
"""[Supply Discovery] Step 2: Content-based file classifier with ownership detection."""
import json,re,subprocess,os
from pathlib import Path
from typing import List,Dict,Tuple,NewType
from pydantic import BaseModel
import pymupdf
from docx import Document
from pptx import Presentation

# Type aliases
FilePath = NewType('FilePath', Path)
Cat      = NewType('Cat', str)      # Category
Reason   = NewType('Reason', str)
Txt      = NewType('Txt', str)      # TextContent

class ClassifiedFiles(BaseModel):
    """Classified file categories."""
    class Config: frozen = True
    user_resumes: List[Dict]
    user_cls: List[Dict]       # cover_letters
    user_combined: List[Dict]
    other_resumes: List[Dict]
    presentations: List[Dict]  # NEW
    jds: List[Dict]            # NEW: Job descriptions
    recruiters: List[Dict]     # NEW: Recruiter info
    companies: List[Dict]      # NEW: Company research
    tracking: List[Dict]       # tracking_files
    other: List[Dict]

# --- Folder-based classification (HIGH PRIORITY) ---
FOLDER_RULES = {
    'Presentations':  'presentation',     # Files in Presentations/ folder
    'Roles':          'jd',               # Job descriptions
    'Jobs':           'jd',               # Job descriptions
    'Recruiters':     'recruiter',        # Recruiter info
    'CV samples':     'other_resume',     # Other people's resumes
    'OthersInfo':     'other_resume',     # Other people's info
}

def get_folder_hint(fp: Path) -> str:
    """Get classification hint from folder path."""
    parts = fp.parts
    for folder, cat in FOLDER_RULES.items():
        if folder in parts:
            return cat
    return None

# --- Text Extraction (one-liners where possible) ---
def extract_pdf(fp: Path) -> Txt:
    try:
        doc = pymupdf.open(fp)
        txt = "".join(p.get_text() for p in doc)
        doc.close()
        return txt.lower()
    except Exception as e:
        print(f"  Warning: Could not read PDF {fp.name}: {e}")
        return ""
def extract_docx(fp: Path) -> Txt:
    try:
        doc = Document(fp)
        return "\n".join(p.text for p in doc.paragraphs).lower()
    except Exception as e:
        print(f"  Warning: Could not read DOCX {fp.name}: {e}")
        return ""
def extract_doc(fp: Path) -> Txt:
    """Extract from old .doc using textutil (Mac) or raw bytes."""
    try:
        r = subprocess.run(['textutil','-convert','txt','-stdout',str(fp)],
                          capture_output=True, text=True, timeout=10)
        if r.returncode == 0: return r.stdout.lower()
        # Fallback: extract ASCII
        with open(fp,'rb') as f:
            raw = f.read()
            return ''.join(chr(b) if 32<=b<127 else ' ' for b in raw).lower()
    except Exception as e:
        print(f"  Warning: Could not read .doc {fp.name}: {e}")
        return ""
def extract_pptx(fp: Path) -> Txt:
    try:
        prs = Presentation(fp)
        return "\n".join(sh.text for sl in prs.slides for sh in sl.shapes if hasattr(sh,"text")).lower()
    except Exception as e:
        print(f"  Warning: Could not read PPTX {fp.name}: {e}")
        return ""
def extract_txt(fp: Path) -> Txt:
    try: return fp.read_text(encoding='utf-8').lower()
    except: return ""

# --- Classification Helpers ---
def is_user_doc(
    txt: str,           # Document text (lowercase)
    fname: str,         # Filename
    name: str=None,     # User's full name
    email: str=None     # User's email
) -> bool:
    """Check if doc belongs to user (name in content/filename OR email in content)."""
    name  = name  or os.getenv('USER_NAME', 'YOUR_NAME')
    email = email or os.getenv('USER_EMAIL')
    nm_lower   = name.lower()
    nm_compact = nm_lower.replace(' ', '')
    fn_lower   = fname.lower().replace(' ', '')
    # Check name in filename/content, email in content
    has_name = nm_compact in fn_lower or nm_lower in txt
    has_email = email and email.lower() in txt
    return has_name or has_email

CL_PATTERNS = [
    r'dear\s+(?:hiring|recruiter|manager|team)',
    r'i\s+am\s+writing\s+to\s+express',
    r'i\s+am\s+(?:excited|pleased|thrilled)\s+to\s+apply',
    r'sincerely', r'best\s+regards',
    r'thank\s+you\s+for\s+(?:considering|your\s+time)',
    r'i\s+would\s+welcome\s+the\s+opportunity',
    r'i\s+look\s+forward\s+to', r'cover\s+letter',
    r'enthusiastic\s+about.*opportunity', r'interest\s+in.*position'
]
def is_cl(txt: str) -> bool:
    """Check if doc is a cover letter."""
    return sum(1 for p in CL_PATTERNS if re.search(p, txt)) >= 1

RES_SECTIONS = [
    r'\bexperience\b', r'\beducation\b', r'\bskills\b',
    r'professional\s+summary', r'accomplishments', r'\bwork\s+history\b',
    r'\bemployment\b', r'technical\s+(?:skills|arsenal|expertise)', r'core\s+capabilities'
]
RES_PATTERNS = {
    'dates':    r'(?:19|20)\d{2}\s*[-–—]\s*(?:(?:19|20)\d{2}|present|current)',
    'titles':   r'\b(?:vp|vice president|director|manager|chief|lead|senior|sr|engineer|scientist|analyst|specialist)\b',
    'companies': r'\b(?:at\s+\w+|worked\s+at|employed\s+by)\b',
    'metrics':   r'(?:\$[\d,]+[kmb]?|\d+%|\d+\+\s+(?:years|people|projects))'
}
def is_resume(txt: str) -> bool:
    """Check if doc has resume structure (3+ sections AND 2+ patterns)."""
    sec_cnt = sum(1 for p in RES_SECTIONS if re.search(p, txt))
    has_bullets = txt.count('•') + txt.count('- ') + txt.count('* ') > 5
    pat_cnt = sum(bool(re.search(v,txt)) for v in RES_PATTERNS.values()) + has_bullets
    return sec_cnt >= 3 and pat_cnt >= 2

JD_PATTERNS = [
    r'\bjob\s+description\b', r'\bresponsibilities\b', r'\brequirements\b',
    r'\bqualifications\b', r'\bwe\s+are\s+(?:looking|seeking)\b',
    r'\bthe\s+ideal\s+candidate\b', r'\breports\s+to\b', r'\bequal\s+opportunity\b'
]
def is_jd(txt: str, fname: str) -> bool:
    """Check if doc is a job description."""
    fn_lower = fname.lower()
    # Filename hints
    if 'job description' in fn_lower or 'jd' in fn_lower: return True
    if 'vice-president' in fn_lower and '--' in fn_lower: return True  # Scraped JD naming
    # Content patterns
    return sum(1 for p in JD_PATTERNS if re.search(p, txt)) >= 3

def is_pres(txt: str, fname: str, fp: Path) -> bool:
    """Check if doc is a presentation (strict: in Presentations folder or specific files)."""
    fn = fname.lower().replace(' ','')
    # User's profile/resume pptx files are NOT presentations
    if 'nitinverma' in fn and fn.endswith('.pptx'):
        return False
    # In Presentations folder
    if 'Presentations' in fp.parts:
        return True
    # Specific presentation files by name
    pres_names = ['enterpriseagenticai', 'hexagonphysicalai']
    if any(p in fn for p in pres_names):
        return True
    return False

def is_company_research(fname: str) -> bool:
    """Check if file is company research/list."""
    fn = fname.lower()
    hints = ['companies', 'company', 'top10k', 'top50k', 'contacts',
             'computerprogramming', 'software_computer', 'hitech', 'jobsearchresults']
    return any(h in fn for h in hints)

def is_recruiter_info(fname: str) -> bool:
    """Check if file is recruiter/exec search info."""
    fn = fname.lower()
    # Must have 'exec' and 'search' together, or 'firms', or 'recruiter'
    has_exec_search = 'exec' in fn and 'search' in fn
    hints = ['firms', 'recruiter']
    return has_exec_search or any(h in fn for h in hints)

# Load other people's names from config (git-ignored)
def _load_other_names() -> List[str]:
    cfg_path = Path(__file__).parent.parent / "config" / "classification_config.json"
    if cfg_path.exists():
        try:
            with open(cfg_path, 'r') as f:
                cfg = json.load(f)
                return cfg.get('other_peoples_names', [])
        except: pass
    return []  # Empty list if no config

OTHER_NAMES = _load_other_names()
def is_other_resume(txt: str) -> bool:
    return any(re.search(p, txt) for p in OTHER_NAMES)

# --- Main Classification ---
def classify_file(finfo: Dict, name: str=None, email: str=None) -> Tuple[Cat, Reason]:
    """Classify a file by analyzing content and folder context."""
    # Skip directories
    if finfo.get('is_dir') or finfo.get('is_directory'):
        return 'skip', 'directory'
    
    fp   = Path(finfo['path'])
    fn   = finfo['name']
    ext  = finfo.get('ext') or finfo.get('extension', '')
    
    # Skip directories (some may be listed without is_dir flag)
    if not ext and fp.is_dir():
        return 'skip', 'directory'
    
    # 1. FOLDER-BASED CLASSIFICATION (highest priority)
    folder_hint = get_folder_hint(fp)
    if folder_hint == 'presentation':
        return 'presentation', 'in Presentations folder'
    if folder_hint == 'jd':
        return 'jd', 'in Roles/Jobs folder (job description)'
    if folder_hint == 'recruiter':
        return 'recruiter', 'in Recruiters folder'
    if folder_hint == 'other_resume':
        return 'other_resume', 'in CV samples/OthersInfo folder'
    
    # 2. LINKEDIN EXPORT DATA (separate category - not useful for resume building)
    if 'Basic_LinkedInDataExport' in str(fp):
        return 'other', 'LinkedIn export data'
    
    # 3. FILENAME-BASED CLASSIFICATION
    fn_lower = fn.lower()
    
    # Job search logs -> JDs category
    if 'work search' in fn_lower or 'job search' in fn_lower:
        return 'jd', 'job search activity log'
    
    if is_company_research(fn):
        return 'company', 'company research/list filename'
    if is_recruiter_info(fn):
        return 'recruiter', 'recruiter/exec search filename'
    
    # 4. SPREADSHEETS
    if ext in ['.xlsx', '.csv']:
        return 'other', 'spreadsheet'
    
    # 4. UNSUPPORTED TYPES
    if ext not in ['.pdf', '.docx', '.doc', '.pptx', '.txt']:
        return 'other', f'unsupported file type {ext}'
    
    # 5. EXTRACT TEXT
    if   ext == '.pdf':  txt = extract_pdf(fp)
    elif ext == '.docx': txt = extract_docx(fp)
    elif ext == '.doc':  txt = extract_doc(fp)
    elif ext == '.pptx': txt = extract_pptx(fp)
    else:                txt = extract_txt(fp)
    
    if not txt: return 'other', 'could not extract text'
    
    # 6. CONTENT-BASED: JD detection
    if is_jd(txt, fn):
        return 'jd', 'job description content'
    
    # 7. CONTENT-BASED: Presentation detection  
    if is_pres(txt, fn, fp):
        return 'presentation', 'presentation content'
    
    # 8. CONTENT-BASED: Other person's resume
    if is_other_resume(txt):
        return 'other_resume', "identified as another person's resume"
    
    # 9. OWNERSHIP CHECK
    if not is_user_doc(txt, fn, name, email):
        if is_resume(txt): return 'other_resume', "resume structure but not user's"
        return 'other', "not identified as user's document"
    
    # 10. USER'S DOC - determine type
    has_cl, has_res = is_cl(txt), is_resume(txt)
    wc = len(txt.split())
    
    if has_cl and has_res:
        if wc < 600: return 'user_cl', "user's cover letter (too short to be combined)"
        return 'user_combined', "user's document with both resume and cover letter"
    if has_cl:  return 'user_cl', "user's cover letter only"
    if has_res: return 'user_resume', "user's resume only"
    return 'user_resume', "user's document (structure unclear, defaulting to resume)"

def classify_inventory(
    inv_path: Path,     # Path to file_inventory.json
    name: str=None,     # User name
    email: str=None     # User email
) -> ClassifiedFiles:
    """Classify all files from inventory."""
    with open(inv_path, 'r', encoding='utf-8') as f:
        inv = json.load(f)
    
    cls = {'user_resumes':[], 'user_cls':[], 'user_combined':[],
           'other_resumes':[], 'presentations':[], 'jds':[], 
           'recruiters':[], 'companies':[], 'tracking':[], 'other':[]}
    
    files = inv['files']
    print(f"Analyzing {len(files)} files...\n")
    
    for i, finfo in enumerate(files, 1):
        if i % 20 == 0: print(f"  Processed {i}/{len(files)} files...")
        
        cat, reason = classify_file(finfo, name, email)
        
        # Skip directories entirely
        if cat == 'skip': continue
        
        finfo['classification_reason'] = reason
        
        # Categorization
        if   cat == 'user_resume':   cls['user_resumes'].append(finfo)
        elif cat == 'user_cl':       cls['user_cls'].append(finfo)
        elif cat == 'user_combined': cls['user_combined'].append(finfo)
        elif cat == 'other_resume':  cls['other_resumes'].append(finfo)
        elif cat == 'presentation':  cls['presentations'].append(finfo)
        elif cat == 'jd':            cls['jds'].append(finfo)
        elif cat == 'recruiter':     cls['recruiters'].append(finfo)
        elif cat == 'company':       cls['companies'].append(finfo)
        elif cat == 'tracking':      cls['tracking'].append(finfo)
        else:                        cls['other'].append(finfo)
    
    return ClassifiedFiles(**cls)

def main(
    inv: Path = Path(__file__).parent.parent/"data"/"supply"/"1_file_inventory.json",
    out: Path = Path(__file__).parent.parent/"data"/"supply"/"2_file_inventory.json",
    name: str = None,   # User's full name (default: USER_NAME env var)
    email: str = None   # User's email (default: USER_EMAIL env var)
):
    """Main execution."""
    name  = os.getenv('USER_NAME', name)
    email = os.getenv('USER_EMAIL', email)
    
    if not inv.exists():
        print(f"Error: File inventory not found at {inv}")
        print("Please run 1_scan_resume_folder.py first.")
        return
    
    print("Classifying files with FOLDER + CONTENT ANALYSIS...")
    print("Using: folder context, filename patterns, content analysis\n")
    
    classified = classify_inventory(inv, name, email)
    
    with open(out, 'w', encoding='utf-8') as f:
        json.dump(classified.model_dump(), f, indent=2, ensure_ascii=False)
    
    print(f"\n✓ Classification complete!")
    print(f"✓ Saved to: {out}\n")
    
    print("Classification Summary:")
    print(f"  User Resumes: {len(classified.user_resumes)} files")
    print(f"  User Cover Letters: {len(classified.user_cls)} files")
    print(f"  User Combined: {len(classified.user_combined)} files")
    print(f"  Presentations: {len(classified.presentations)} files")
    print(f"  Job Descriptions: {len(classified.jds)} files")
    print(f"  Other Resumes: {len(classified.other_resumes)} files")
    print(f"  Recruiters: {len(classified.recruiters)} files")
    print(f"  Companies: {len(classified.companies)} files")
    print(f"  Tracking Files: {len(classified.tracking)} files")
    print(f"  Other Files: {len(classified.other)} files")
    
    total_user = len(classified.user_resumes) + len(classified.user_cls) + len(classified.user_combined)
    print(f"\n📊 Total User Documents: {total_user}")

if __name__ == "__main__": main()
